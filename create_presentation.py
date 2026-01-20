#!/usr/bin/env python3
"""
Create a professional 4-slide PowerPoint presentation for Ghost Protocol
Tailored for Envoy Ortus as a major client
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme - Professional, sophisticated
PRIMARY_COLOR = RGBColor(15, 23, 42)  # Deep navy
ACCENT_COLOR = RGBColor(59, 130, 246)  # Professional blue
TEXT_COLOR = RGBColor(51, 65, 85)  # Dark gray
LIGHT_BG = RGBColor(248, 250, 252)  # Very light blue
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR
    
    # Add title in header
    title_frame = header.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(12)
    title_p.space_after = Pt(0)
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    
    # Add accent line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.color.rgb = ACCENT_COLOR
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        
        # Add bullet
        p.font.bold = False
    
    return slide

# Slide 1: Title/Introduction
add_title_slide(
    prs,
    "Ghost Protocol",
    "Enterprise-Grade Digital Solutions for Strategic Transformation"
)

# Slide 2: Company Introduction & Core Values
slide2_content = [
    "• Established presence in premium IT solutions and digital transformation",
    "• Specialized expertise in enterprise-scale web development and custom software engineering",
    "• Trusted partner for mission-critical infrastructure and digital infrastructure optimization",
    "",
    "Core Values:",
    "• Technical Excellence: Unwavering commitment to architecture and code quality",
    "• Strategic Partnership: Long-term client success over transactional relationships",
    "• Operational Integrity: Transparent, accountable, and results-driven delivery"
]
add_content_slide(prs, "About Ghost Protocol", slide2_content)

# Slide 3: Products & Services
slide3_content = [
    "Solutions Tailored for Enterprise Environments:",
    "",
    "• Web Development: Full-stack applications, scalable architectures, cloud-native solutions",
    "",
    "• Software Engineering: Custom applications, API design, system integration, legacy modernization",
    "",
    "• Digital Infrastructure: Cloud architecture optimization, DevOps strategy, security hardening",
    "",
    "• Strategic Consulting: Technology roadmapping, architecture design, digital transformation planning"
]
add_content_slide(prs, "Products & Services", slide3_content)

# Slide 4: Key Benefits for Envoy Ortus
slide4_content = [
    "Differentiated Value Proposition:",
    "",
    "• Specialized Expertise: Deep technical knowledge in enterprise-grade systems and complex requirements",
    "",
    "• Operational Excellence: Proven ability to deliver high-quality solutions on schedule and within specifications",
    "",
    "• Long-Term Partnership: Committed to understanding your strategic objectives and contributing to sustained success",
    "",
    "• Scalability & Future-Proofing: Solutions architected for growth, supporting your evolution and market expansion"
]
add_content_slide(prs, "Strategic Value", slide4_content)

# Slide 5: Pricing Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Add header bar
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY_COLOR
header.line.color.rgb = PRIMARY_COLOR

title_frame = header.text_frame
title_frame.clear()
title_p = title_frame.paragraphs[0]
title_p.text = "Investment & Engagement Models"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = WHITE
title_p.space_before = Pt(12)
title_p.space_after = Pt(0)
title_frame.margin_left = Inches(0.5)
title_frame.margin_right = Inches(0.5)

# Add accent line
line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_COLOR
line.line.color.rgb = ACCENT_COLOR

# Add content
content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.2))
text_frame = content_box.text_frame
text_frame.word_wrap = True

pricing_content = [
    "Flexible engagement models tailored to project scope and timeline:",
    "",
    "• Project-Based Engagements",
    "  Fixed scope, deliverables, and timeline with transparent cost structure",
    "",
    "• Time & Materials (T&M)",
    "  Flexible resource allocation for evolving requirements",
    "",
    "• Retainer Partnerships",
    "  Ongoing strategic support and optimization services",
    "",
    "Detailed proposal and pricing structure to be prepared by sales team"
]

for i, item in enumerate(pricing_content):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(16) if item.startswith("•") and not item.startswith("  ") else Pt(15)
    p.font.color.rgb = TEXT_COLOR
    
    if item.startswith("  "):
        p.level = 1
        p.font.size = Pt(15)
    
    p.space_before = Pt(4)
    p.space_after = Pt(4)

# Save presentation
output_path = "/home/kami/Git Projects/ghosts-lk.github.io/Ghost_Protocol_Envoy_Ortus_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"   Slides: 5")
print(f"   Format: Professional, clean design")
print(f"   Ready for client presentation")
